import PyPDF2
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
from pdf2image import convert_from_path
from pptx import Presentation
import json
import os
import io
import cv2
import numpy as np
from typing import Dict, List, Any, Optional
import re
from datetime import datetime

class DocumentExtractor:
    def __init__(self):
        self.extracted_data = {}
    
    def extract_pdf_text(self, pdf_path: str, source_file_type: str = "unknown") -> Dict[str, Any]:
        """Enhanced PDF text extraction with source file tracking and better error handling"""
        try:
            # Open document with error handling
            try:
                doc = fitz.open(pdf_path)
            except Exception as e:
                print(f"         ❌ Cannot open PDF: {e}")
                return self._create_error_result(pdf_path, source_file_type, f"Cannot open PDF: {e}")
            
            full_text = ""
            images_text = ""
            structured_content = []
            
            # Get page count
            page_count = len(doc)
            if page_count == 0:
                print(f"         ❌ PDF has no pages")
                doc.close()
                return self._create_error_result(pdf_path, source_file_type, "PDF has no pages")
            
            print(f"         📖 Processing {page_count} pages...")
            
            for page_num in range(page_count):
                try:
                    page = doc.load_page(page_num)
                    
                    # ENHANCED: Extract text with better formatting preservation
                    text = page.get_text()
                    
                    # IMPROVED: Also try to get text with layout information
                    try:
                        text_dict = page.get_text("dict")
                        layout_text = self._extract_text_with_layout(text_dict)
                        if layout_text and len(layout_text) > len(text):
                            text = layout_text
                    except Exception as layout_error:
                        print(f"            ⚠️ Layout extraction failed for page {page_num + 1}: {layout_error}")
                    
                    if text.strip():  # Only add non-empty pages
                        # ENHANCED: Include source file type in page header
                        page_header = f"\n{'='*50}\nSOURCE_FILE: {source_file_type.upper()}\nPAGE {page_num + 1}\n{'='*50}\n"
                        
                        # IMPROVED: Better text cleaning and preservation
                        cleaned_text = self._clean_and_preserve_text(text)
                        full_text += page_header + cleaned_text + "\n"
                        
                        # Try to structure the content better
                        structured_content.append({
                            'source_file': source_file_type,
                            'page': page_num + 1,
                            'content': cleaned_text,
                            'length': len(cleaned_text)
                        })
                    
                    # Extract images and perform OCR with better error handling
                    try:
                        image_list = page.get_images()
                        if image_list:
                            print(f"            🖼️  Found {len(image_list)} images on page {page_num + 1}")
                        
                        for img_index, img in enumerate(image_list):
                            try:
                                xref = img[0]
                                pix = fitz.Pixmap(doc, xref)
                                
                                # Check if it's a suitable image for OCR
                                if pix.n - pix.alpha < 4:  # GRAY or RGB
                                    img_data = pix.tobytes("png")
                                    img_pil = Image.open(io.BytesIO(img_data))
                                    
                                    # Perform OCR on the image with timeout protection
                                    try:
                                        ocr_text = pytesseract.image_to_string(img_pil, config='--psm 6')
                                        if ocr_text.strip():  # Only add if OCR found text
                                            # ENHANCED: Include source file in OCR header
                                            ocr_header = f"\n[OCR from {source_file_type.upper()}: Page {page_num + 1}, Image {img_index + 1}]\n"
                                            images_text += ocr_header + ocr_text + "\n"
                                    except Exception as ocr_error:
                                        print(f"            ⚠️ OCR failed for image {img_index + 1}: {ocr_error}")
                                    
                                    # Clean up PIL image
                                    img_pil.close()
                                
                                # Clean up pixmap immediately
                                pix = None
                                
                            except Exception as img_error:
                                print(f"            ⚠️ Error processing image {img_index + 1} on page {page_num + 1}: {img_error}")
                                continue
                    except Exception as images_error:
                        print(f"            ⚠️ Error processing images on page {page_num + 1}: {images_error}")
                    
                    # Clean up page reference
                    page = None
                    
                except Exception as page_error:
                    print(f"         ⚠️ Error processing page {page_num + 1}: {page_error}")
                    continue
            
            # Close document properly
            doc.close()
            
            # Validate extracted content
            if not full_text.strip() and not images_text.strip():
                print(f"         ⚠️ No text content extracted from PDF")
                return self._create_error_result(pdf_path, source_file_type, "No text content extracted")
            
            # Create enhanced combined text with better structure
            combined_text = self.create_enhanced_combined_text(full_text, images_text, structured_content)
            
            print(f"         ✅ Extracted {len(full_text):,} chars text + {len(images_text):,} chars OCR")
            
            return {
                "full_text": full_text,
                "images_text": images_text,
                "combined_text": combined_text,
                "structured_content": structured_content,
                "file_path": pdf_path,
                "source_file_type": source_file_type,
                "pages": page_count,
                "text_length": len(full_text),
                "ocr_length": len(images_text),
                "total_length": len(combined_text),
                "extraction_success": True
            }
            
        except Exception as e:
            print(f"         ❌ Error extracting PDF text from {os.path.basename(pdf_path)}: {e}")
            return self._create_error_result(pdf_path, source_file_type, str(e))
    
    def _extract_text_with_layout(self, text_dict: Dict) -> str:
        """IMPROVED: Extract text while preserving layout and formatting"""
        try:
            extracted_text = ""
            
            if "blocks" in text_dict:
                for block in text_dict["blocks"]:
                    if "lines" in block:
                        block_text = ""
                        for line in block["lines"]:
                            if "spans" in line:
                                line_text = ""
                                for span in line["spans"]:
                                    if "text" in span:
                                        text = span["text"]
                                        # Preserve formatting information
                                        if span.get("flags", 0) & 2**4:  # Bold
                                            text = f"**{text}**"
                                        line_text += text
                                if line_text.strip():
                                    block_text += line_text + "\n"
                        if block_text.strip():
                            extracted_text += block_text + "\n"
            
            return extracted_text
        except Exception as e:
            print(f"            ⚠️ Layout text extraction error: {e}")
            return ""
    
    def _clean_and_preserve_text(self, text: str) -> str:
        """IMPROVED: Clean text while preserving important formatting"""
        if not text:
            return ""
        
        # Split into lines for processing
        lines = text.split('\n')
        cleaned_lines = []
        
        for line in lines:
            # Remove excessive whitespace but preserve structure
            cleaned_line = ' '.join(line.split())
            
            # Preserve bullet points and list items
            if cleaned_line:
                # Detect bullet points and preserve them
                if re.match(r'^\s*[•·▪▫◦‣⁃]\s*', line) or re.match(r'^\s*[-*]\s+', line):
                    cleaned_lines.append(f"• {cleaned_line.lstrip('•·▪▫◦‣⁃-* ')}")
                elif re.match(r'^\s*\d+\.\s+', line):
                    cleaned_lines.append(cleaned_line)
                else:
                    cleaned_lines.append(cleaned_line)
        
        return '\n'.join(cleaned_lines)
    
    def _create_error_result(self, file_path: str, source_file_type: str, error_msg: str) -> Dict[str, Any]:
        """Create standardized error result"""
        return {
            "full_text": "", 
            "images_text": "", 
            "combined_text": "", 
            "structured_content": [],
            "file_path": file_path,
            "source_file_type": source_file_type,
            "pages": 0,
            "text_length": 0,
            "ocr_length": 0,
            "total_length": 0,
            "extraction_success": False,
            "error": error_msg
        }
    
    def create_enhanced_combined_text(self, full_text: str, images_text: str, structured_content: List[Dict]) -> str:
        """Create enhanced combined text with better structure for semantic matching"""
        
        # Start with key information extraction
        enhanced_text = "KEY INFORMATION EXTRACTED:\n" + "="*60 + "\n"
        
        # Extract key information patterns
        key_info = self.extract_key_information(full_text)
        for category, info in key_info.items():
            if info:
                enhanced_text += f"\n{category.upper()}:\n"
                for item in info[:5]:  # Limit to 5 items per category for token efficiency
                    enhanced_text += f"- {item}\n"
        
        enhanced_text += "\n" + "="*60 + "\nFULL DOCUMENT CONTENT:\n" + "="*60 + "\n"
        enhanced_text += full_text
        
        if images_text.strip():
            enhanced_text += "\n" + "="*60 + "\nOCR EXTRACTED CONTENT:\n" + "="*60 + "\n"
            enhanced_text += images_text
        
        return enhanced_text
    
    def extract_key_information(self, text: str) -> Dict[str, List[str]]:
        """Extract key information patterns from text for better semantic matching"""
        key_info = {
            'company_names': [],
            'financial_information': [],
            'people_and_roles': [],
            'locations': [],
            'dates_and_years': [],
            'industries_and_verticals': [],
            'funding_and_investment': []
        }
        
        if not text:
            return key_info
        
        try:
            # Company name patterns - improved
            company_patterns = [
                r'(?:Company|Business|Startup|Firm|Corporation)[:\s]+([^\n,.]{3,50})',
                r'([A-Z][a-zA-Z]+(?:\s+[A-Z][a-zA-Z]*)*)\s*[|\-|•]\s*(?:Private|Company|Profile)',
                r'(?:^|\n)([A-Z][a-zA-Z]+(?:\s+[A-Z][a-zA-Z]*)*)\s*(?:\||Private Company Profile)',
                r'(?:^|\n)([A-Z][a-zA-Z]+(?:\s+[A-Z][a-zA-Z]*)*)\s*(?:Inc\.|LLC|Corp\.|Ltd\.)'
            ]
            
            for pattern in company_patterns:
                try:
                    matches = re.finditer(pattern, text, re.MULTILINE | re.IGNORECASE)
                    for match in matches:
                        company_name = match.group(1).strip()
                        if len(company_name) > 2 and company_name not in key_info['company_names']:
                            key_info['company_names'].append(company_name)
                except Exception as e:
                    print(f"            ⚠️ Error in company pattern matching: {e}")
                    continue
            
            # Financial information - improved
            financial_patterns = [
                r'\$[\d,]+(?:\.\d+)?[MBK]?(?:\s*(?:million|billion|thousand))?',
                r'(?:valuation|raised|funding|investment|revenue)[:\s]*\$[\d,]+(?:\.\d+)?[MBK]?',
                r'(?:Series [A-Z]|Seed|Pre-seed|Growth)(?:\s+(?:funding|round|investment))?',
                r'(?:pre-money|post-money)[:\s]*\$[\d,]+(?:\.\d+)?[MBK]?',
                r'(?:ARR|MRR|revenue)[:\s]*\$[\d,]+(?:\.\d+)?[MBK]?'
            ]
            
            for pattern in financial_patterns:
                try:
                    matches = re.finditer(pattern, text, re.IGNORECASE)
                    for match in matches:
                        financial_info = match.group().strip()
                        if financial_info not in key_info['financial_information']:
                            key_info['financial_information'].append(financial_info)
                except Exception as e:
                    print(f"            ⚠️ Error in financial pattern matching: {e}")
                    continue
            
            # People and roles - improved
            people_patterns = [
                r'(?:CEO|CTO|CFO|Founder|Co-Founder|Chief|President|Director)[:\s]*([A-Z][a-zA-Z\s]+)',
                r'([A-Z][a-zA-Z]+\s+[A-Z][a-zA-Z]+)(?:\s*-?\s*(?:CEO|CTO|CFO|Founder|Co-Founder))',
                r'(?:Founded by|Led by|Team)[:\s]*([A-Z][a-zA-Z\s,]+)',
                r'([A-Z][a-zA-Z]+\s+[A-Z][a-zA-Z]+),?\s*(?:CEO|CTO|CFO|Founder)'
            ]
            
            for pattern in people_patterns:
                try:
                    matches = re.finditer(pattern, text, re.IGNORECASE)
                    for match in matches:
                        person_info = match.group().strip()
                        if len(person_info) > 5 and person_info not in key_info['people_and_roles']:
                            key_info['people_and_roles'].append(person_info)
                except Exception as e:
                    print(f"            ⚠️ Error in people pattern matching: {e}")
                    continue
            
            # Locations - improved
            location_patterns = [
                r'(?:located|based|headquarters|office|address)[:\s]*([A-Z][a-zA-Z\s,]+(?:CA|NY|TX|FL|USA|United States))',
                r'([A-Z][a-zA-Z]+,\s*[A-Z]{2}(?:,\s*(?:USA|United States))?)',
                r'(?:San Francisco|New York|Los Angeles|Boston|Seattle|Austin|Chicago|Miami|Palo Alto),?\s*[A-Z]{2}?',
                r'(?:California|New York|Texas|Florida|Massachusetts|Washington|Illinois),?\s*USA?'
            ]
            
            for pattern in location_patterns:
                try:
                    matches = re.finditer(pattern, text, re.IGNORECASE)
                    for match in matches:
                        location = match.group().strip()
                        if len(location) > 3 and location not in key_info['locations']:
                            key_info['locations'].append(location)
                except Exception as e:
                    print(f"            ⚠️ Error in location pattern matching: {e}")
                    continue
            
            # Years and dates - improved
            year_patterns = [
                r'(?:founded|established|started|incorporated|year)[:\s]*(\d{4})',
                r'(\d{4})(?:\s*-\s*(?:founded|established|started))',
                r'(?:as of|dated?)[:\s]*(\d{1,2}[-/]\d{1,2}[-/]\d{4})',
                r'(?:since|from)\s*(\d{4})'
            ]
            
            for pattern in year_patterns:
                try:
                    matches = re.finditer(pattern, text, re.IGNORECASE)
                    for match in matches:
                        date_info = match.group().strip()
                        if date_info not in key_info['dates_and_years']:
                            key_info['dates_and_years'].append(date_info)
                except Exception as e:
                    print(f"            ⚠️ Error in date pattern matching: {e}")
                    continue
            
            # Industries and verticals - improved
            industry_patterns = [
                r'(?:industry|sector|market|vertical)[:\s]*([A-Z][a-zA-Z\s&,]+)',
                r'(?:AI|Artificial Intelligence|Machine Learning|Software|Technology|Healthcare|Fintech|Biotech|SaaS)',
                r'(?:Media and Information Services|Business Software|Enterprise|B2B|B2C)',
                r'(?:EdTech|PropTech|HealthTech|CleanTech|FoodTech|RetailTech)'
            ]
            
            for pattern in industry_patterns:
                try:
                    matches = re.finditer(pattern, text, re.IGNORECASE)
                    for match in matches:
                        industry = match.group().strip()
                        if len(industry) > 2 and industry not in key_info['industries_and_verticals']:
                            key_info['industries_and_verticals'].append(industry)
                except Exception as e:
                    print(f"            ⚠️ Error in industry pattern matching: {e}")
                    continue
            
        except Exception as e:
            print(f"            ⚠️ Error in key information extraction: {e}")
        
        return key_info
    
    def extract_pptx_text(self, pptx_path: str, source_file_type: str = "unknown") -> Dict[str, Any]:
        """Enhanced PPTX text extraction with source file tracking and better error handling"""
        try:
            # Open presentation with error handling
            try:
                prs = Presentation(pptx_path)
            except Exception as e:
                print(f"         ❌ Cannot open PPTX: {e}")
                return self._create_error_result(pptx_path, source_file_type, f"Cannot open PPTX: {e}")
            
            full_text = ""
            images_text = ""
            slide_count = len(prs.slides)
            structured_content = []
            
            if slide_count == 0:
                print(f"         ❌ PPTX has no slides")
                return self._create_error_result(pptx_path, source_file_type, "PPTX has no slides")
            
            print(f"         📊 Processing {slide_count} slides...")
            
            for slide_num, slide in enumerate(prs.slides):
                try:
                    # ENHANCED: Add source file info to slide separator
                    slide_header = f"\n{'='*40}\nSOURCE_FILE: {source_file_type.upper()}\nSLIDE {slide_num + 1}\n{'='*40}\n"
                    slide_text = ""
                    
                    # Extract text from shapes with better error handling
                    shape_count = 0
                    for shape in slide.shapes:
                        try:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text += shape.text + "\n"
                                shape_count += 1
                            
                            # Handle images (basic detection)
                            if hasattr(shape, 'shape_type') and shape.shape_type == 13:  # Picture
                                try:
                                    # ENHANCED: Include source file in image detection
                                    images_text += f"[IMAGE DETECTED in {source_file_type.upper()}: Slide {slide_num + 1}]\n"
                                    print(f"            🖼️  Image detected in slide {slide_num + 1}")
                                except Exception as e:
                                    print(f"            ⚠️ Error processing image in slide {slide_num + 1}: {e}")
                        except Exception as shape_error:
                            print(f"            ⚠️ Error processing shape in slide {slide_num + 1}: {shape_error}")
                            continue
                    
                    if slide_text.strip():
                        full_text += slide_header + slide_text
                        structured_content.append({
                            'source_file': source_file_type,
                            'slide': slide_num + 1,
                            'content': slide_text,
                            'length': len(slide_text)
                        })
                        
                        if shape_count > 0:
                            print(f"            📝 Extracted text from {shape_count} shapes")
                            
                except Exception as slide_error:
                    print(f"         ⚠️ Error processing slide {slide_num + 1}: {slide_error}")
                    continue
            
            # Validate extracted content
            if not full_text.strip():
                print(f"         ⚠️ No text content extracted from PPTX")
                return self._create_error_result(pptx_path, source_file_type, "No text content extracted")
            
            # Create enhanced combined text
            combined_text = self.create_enhanced_combined_text(full_text, images_text, structured_content)
            
            print(f"         ✅ Extracted {len(full_text):,} chars from {slide_count} slides")
            
            return {
                "full_text": full_text,
                "images_text": images_text,
                "combined_text": combined_text,
                "structured_content": structured_content,
                "file_path": pptx_path,
                "source_file_type": source_file_type,
                "slides": slide_count,
                "text_length": len(full_text),
                "images_detected": images_text.count("[IMAGE DETECTED"),
                "total_length": len(combined_text),
                "extraction_success": True
            }
            
        except Exception as e:
            print(f"         ❌ Error extracting PPTX text from {os.path.basename(pptx_path)}: {e}")
            return self._create_error_result(pptx_path, source_file_type, str(e))
    
    def extract_memo_sections_from_ai_output(self, memo_text: str) -> Dict[str, str]:
        """COMPLETELY IMPROVED: Enhanced memo section extraction with robust cross-page content detection"""
        
        # EXPANDED: All sections as defined in config
        sections = {
            "executive_summary": "",
            "company_information": "",
            "startup_stage": "",
            "deal_summary": "",
            "management_team": "",
            "key_metrics": "",
            "customer_problem": "",
            "product_service_summary": "",
            "investment_themes": "",
            "market_overview": "",
            "list_of_competitors": "",
            "competitive_advantage_summary": "",
            "investment_considerations_risk_factors": ""
        }
        
        if not memo_text or not memo_text.strip():
            print(f"            ⚠️ Empty memo text provided")
            return sections
        
        print(f"            🔍 IMPROVED: Analyzing memo text for robust section extraction...")
        
        try:
            # IMPROVED: Pre-process the text to handle page breaks and formatting
            processed_text = self._preprocess_memo_text(memo_text)
            
            # ENHANCED: Create comprehensive section patterns with multiple alternatives
            section_patterns = self._create_comprehensive_section_patterns()
            
            # IMPROVED: Advanced section detection with cross-page support
            detected_sections = self._detect_sections_advanced(processed_text, section_patterns)
            
            # ENHANCED: Extract content with better cross-page handling
            extracted_sections = self._extract_section_content_advanced(processed_text, detected_sections)
            
            # Merge with default sections
            for section_key, content in extracted_sections.items():
                if section_key in sections:
                    sections[section_key] = content
            
            # IMPROVED: Validation and reporting
            sections_found = len([s for s in sections.values() if s.strip()])
            print(f"            ✅ IMPROVED: Successfully extracted {sections_found}/{len(sections)} sections")
            
            # ENHANCED: Detailed logging of found sections
            for section_name, content in sections.items():
                if content.strip():
                    print(f"               • {section_name}: {len(content)} characters")
            
            # IMPROVED: Debug missing sections
            missing_sections = [name for name, content in sections.items() if not content.strip()]
            if missing_sections:
                print(f"            ⚠️ Missing sections: {', '.join(missing_sections)}")
                # Try alternative extraction for missing sections
                self._try_alternative_extraction(processed_text, sections, missing_sections)
                
        except Exception as e:
            print(f"            ❌ Error in IMPROVED memo section extraction: {e}")
            import traceback
            traceback.print_exc()
        
        return sections
    
    def _preprocess_memo_text(self, memo_text: str) -> str:
        """IMPROVED: Preprocess memo text to handle page breaks and formatting issues"""
        
        # ENHANCED: Better handling of page breaks and formatting
        lines = memo_text.split('\n')
        processed_lines = []
        
        i = 0
        while i < len(lines):
            line = lines[i].strip()
            
            if not line:
                processed_lines.append('')
                i += 1
                continue
            
            # IMPROVED: Handle section headers that might be split across lines
            if self._is_potential_section_header(line):
                # Check if this might be a section header
                potential_header = line
                
                # Look ahead for continuation
                j = i + 1
                while j < len(lines) and j < i + 3:  # Look ahead max 3 lines
                    next_line = lines[j].strip()
                    if not next_line:
                        break
                    if self._is_header_continuation(line, next_line):
                        potential_header += " " + next_line
                        i = j  # Skip the continuation line
                    break
                
                processed_lines.append(potential_header)
            else:
                processed_lines.append(line)
            
            i += 1
        
        # IMPROVED: Join with better spacing preservation
        processed_text = '\n'.join(processed_lines)
        
        # ENHANCED: Clean up excessive whitespace but preserve structure
        processed_text = re.sub(r'\n\s*\n\s*\n+', '\n\n', processed_text)
        
        return processed_text
    
    def _is_potential_section_header(self, line: str) -> bool:
        """Check if a line could be a section header"""
        if len(line) < 3 or len(line) > 100:
            return False
        
        # Common section header patterns
        header_indicators = [
            'executive', 'summary', 'company', 'information', 'startup', 'stage',
            'deal', 'management', 'team', 'metrics', 'customer', 'problem',
            'product', 'service', 'investment', 'themes', 'market', 'overview',
            'competitors', 'competitive', 'advantage', 'considerations', 'risk', 'factors'
        ]
        
        line_lower = line.lower()
        return any(indicator in line_lower for indicator in header_indicators)
    
    def _is_header_continuation(self, first_line: str, second_line: str) -> bool:
        """Check if second line continues the first line header"""
        # Simple heuristics for header continuation
        if len(second_line) > 50:  # Too long to be header continuation
            return False
        
        combined = f"{first_line} {second_line}".lower()
        known_headers = [
            'executive summary',
            'company information', 
            'startup stage',
            'deal summary',
            'management team',
            'key metrics',
            'customer problem',
            'product and service summary',
            'product service summary',
            'investment themes',
            'market overview',
            'list of competitors',
            'competitive advantage summary',
            'investment considerations risk factors',
            'investment considerations & risk factors'
        ]
        
        return any(header in combined for header in known_headers)
    
    def _create_comprehensive_section_patterns(self) -> Dict[str, List[str]]:
        """IMPROVED: Create comprehensive patterns for all section types"""
        
        return {
            "executive_summary": [
                r"executive\s+summary",
                r"exec\s+summary", 
                r"summary",
                r"overview"
            ],
            "company_information": [
                r"company\s+information",
                r"company\s+overview", 
                r"company\s+description",
                r"company\s+profile",
                r"about\s+the\s+company"
            ],
            "startup_stage": [
                r"startup\s+stage",
                r"funding\s+stage", 
                r"stage",
                r"development\s+stage"
            ],
            "deal_summary": [
                r"deal\s+summary",
                r"investment\s+summary", 
                r"funding\s+summary",
                r"transaction\s+summary"
            ],
            "management_team": [
                r"management\s+team",
                r"leadership\s+team", 
                r"team",
                r"founders",
                r"key\s+personnel"
            ],
            "key_metrics": [
                r"key\s+metrics",
                r"metrics", 
                r"financial\s+metrics",
                r"performance\s+metrics",
                r"kpis?"
            ],
            "customer_problem": [
                r"customer\s+problem",
                r"problem", 
                r"customer\s+challenges",
                r"market\s+problem",
                r"pain\s+point"
            ],
            "product_service_summary": [
                r"product\s+and\s+service\s+summary",
                r"product.*service.*summary", 
                r"product\s+summary",
                r"service\s+summary",
                r"solution",
                r"offering"
            ],
            "investment_themes": [
                r"investment\s+themes",
                r"investment\s+thesis",
                r"thesis"
            ],
            "market_overview": [
                r"market\s+overview",
                r"market\s+analysis", 
                r"industry\s+overview",
                r"market\s+opportunity"
            ],
            "list_of_competitors": [
                r"list\s+of\s+competitors",
                r"competitors?\s+list",
                r"competitive\s+landscape",
                r"competition"
            ],
            "competitive_advantage_summary": [
                r"competitive\s+advantage\s+summary",
                r"competitive\s+advantage",
                r"differentiation"
            ],
            "investment_considerations_risk_factors": [
                r"investment\s+considerations.*risk\s+factors",
                r"investment\s+considerations\s*[\&\+]\s*risk\s+factors",
                r"risk\s+factors",
                r"risks?\s+and\s+considerations?"
            ]
        }
    
    def _detect_sections_advanced(self, text: str, section_patterns: Dict[str, List[str]]) -> Dict[str, Dict]:
        """IMPROVED: Advanced section detection with position tracking"""
        
        detected = {}
        lines = text.split('\n')
        
        print(f"            🔍 Scanning {len(lines)} lines for section headers...")
        
        for line_num, line in enumerate(lines):
            original_line = line.strip()
            line_lower = original_line.lower()
            
            if not original_line or len(original_line) < 3:
                continue
            
            # IMPROVED: Try exact matches first, then pattern matches
            for section_name, patterns in section_patterns.items():
                matched = False
                match_type = "none"
                confidence = 0
                
                # Clean line for matching
                clean_line = re.sub(r'[^\w\s&]', '', line_lower).strip()
                
                # ENHANCED: Exact match priority
                exact_matches = {
                    "executive summary": "executive_summary",
                    "company information": "company_information", 
                    "startup stage": "startup_stage",
                    "deal summary": "deal_summary",
                    "management team": "management_team",
                    "key metrics": "key_metrics",
                    "customer problem": "customer_problem",
                    "product and service summary": "product_service_summary",
                    "product service summary": "product_service_summary",
                    "investment themes": "investment_themes",
                    "market overview": "market_overview",
                    "list of competitors": "list_of_competitors",
                    "competitive advantage summary": "competitive_advantage_summary",
                    "investment considerations & risk factors": "investment_considerations_risk_factors",
                    "investment considerations risk factors": "investment_considerations_risk_factors"
                }
                
                if clean_line in exact_matches and exact_matches[clean_line] == section_name:
                    matched = True
                    match_type = "exact"
                    confidence = 100
                
                # IMPROVED: Pattern matching with confidence scoring
                if not matched:
                    for pattern in patterns:
                        try:
                            if re.search(pattern, line_lower, re.IGNORECASE):
                                matched = True
                                match_type = "pattern"
                                confidence = 80 + len(pattern)  # Longer patterns get higher confidence
                                break
                        except Exception as e:
                            print(f"               ⚠️ Pattern error for {pattern}: {e}")
                            continue
                
                if matched:
                    # IMPROVED: Store detection with metadata
                    if section_name not in detected or detected[section_name]['confidence'] < confidence:
                        detected[section_name] = {
                            'line_num': line_num,
                            'line_text': original_line,
                            'match_type': match_type,
                            'confidence': confidence,
                            'clean_text': clean_line
                        }
                        print(f"               ✓ FOUND: '{section_name}' at line {line_num + 1} ({match_type}, confidence: {confidence})")
        
        print(f"            📊 Detected {len(detected)} sections with advanced method")
        return detected
    
    def _extract_section_content_advanced(self, text: str, detected_sections: Dict[str, Dict]) -> Dict[str, str]:
        """IMPROVED: Advanced content extraction with cross-page support"""
        
        extracted = {}
        lines = text.split('\n')
        
        # IMPROVED: Sort sections by line number for proper processing
        sorted_sections = sorted(detected_sections.items(), key=lambda x: x[1]['line_num'])
        
        for i, (section_name, section_info) in enumerate(sorted_sections):
            start_line = section_info['line_num']
            
            # IMPROVED: Determine end line (next section or end of document)
            if i + 1 < len(sorted_sections):
                end_line = sorted_sections[i + 1][1]['line_num']
            else:
                end_line = len(lines)
            
            # ENHANCED: Extract content with better handling
            content_lines = []
            
            # Start from the line after the section header
            for line_idx in range(start_line + 1, end_line):
                if line_idx >= len(lines):
                    break
                
                line = lines[line_idx].strip()
                
                # IMPROVED: Skip empty lines at the beginning but preserve structure within content
                if not content_lines and not line:
                    continue
                
                # ENHANCED: Stop if we hit another section header (safety check)
                if line and self._is_potential_section_header(line) and line_idx != start_line + 1:
                    # Double-check this isn't part of the content
                    if self._is_definitely_new_section(line, section_name):
                        break
                
                content_lines.append(line)
            
            # IMPROVED: Post-process content
            if content_lines:
                # Remove trailing empty lines
                while content_lines and not content_lines[-1]:
                    content_lines.pop()
                
                content = '\n'.join(content_lines).strip()
                
                # ENHANCED: Quality check - ensure we have substantial content
                if len(content) >= 5:  # Very minimal requirement
                    extracted[section_name] = content
                    print(f"               ✅ Extracted '{section_name}': {len(content)} chars")
                else:
                    print(f"               ⚠️ '{section_name}' too short: {len(content)} chars")
            else:
                print(f"               ⚠️ No content found for '{section_name}'")
        
        return extracted
    
    def _is_definitely_new_section(self, line: str, current_section: str) -> bool:
        """Check if line is definitely a new section header"""
        # More sophisticated check to avoid false positives
        line_lower = line.lower()
        
        # Known section starters that are definitely new sections
        definite_sections = [
            'executive summary', 'company information', 'startup stage',
            'deal summary', 'management team', 'key metrics',
            'customer problem', 'product and service summary',
            'investment themes', 'market overview', 'list of competitors',
            'competitive advantage summary', 'investment considerations'
        ]
        
        return any(section in line_lower for section in definite_sections)
    
    def _try_alternative_extraction(self, text: str, sections: Dict[str, str], missing_sections: List[str]):
        """IMPROVED: Try alternative extraction methods for missing sections"""
        
        print(f"            🔄 Trying alternative extraction for {len(missing_sections)} missing sections...")
        
        # Try fuzzy matching and keyword-based extraction
        for section_name in missing_sections:
            content = self._fuzzy_section_extraction(text, section_name)
            if content and len(content.strip()) >= 10:
                sections[section_name] = content
                print(f"               ✅ Alternative extraction successful for '{section_name}': {len(content)} chars")
    
    def _fuzzy_section_extraction(self, text: str, section_name: str) -> str:
        """IMPROVED: Fuzzy extraction based on keywords and context"""
        
        # Define keywords for each section type
        keywords_map = {
            "executive_summary": ["executive", "summary", "overview", "brief"],
            "company_information": ["company", "business", "organization", "firm"],
            "startup_stage": ["stage", "funding", "round", "series"],
            "deal_summary": ["deal", "transaction", "investment", "funding"],
            "management_team": ["team", "management", "leadership", "founders", "ceo", "cto"],
            "key_metrics": ["metrics", "kpi", "performance", "financial", "revenue"],
            "customer_problem": ["problem", "challenge", "pain", "issue", "customer"],
            "product_service_summary": ["product", "service", "solution", "offering"],
            "investment_themes": ["investment", "themes", "thesis", "rationale"],
            "market_overview": ["market", "industry", "sector", "opportunity"],
            "list_of_competitors": ["competitors", "competition", "competitive"],
            "competitive_advantage_summary": ["advantage", "differentiation", "unique"],
            "investment_considerations_risk_factors": ["risk", "considerations", "factors", "challenges"]
        }
        
        if section_name not in keywords_map:
            return ""
        
        keywords = keywords_map[section_name]
        
        # Find paragraphs that contain multiple keywords from this section
        paragraphs = text.split('\n\n')
        best_paragraph = ""
        best_score = 0
        
        for paragraph in paragraphs:
            if len(paragraph.strip()) < 20:
                continue
            
            score = 0
            paragraph_lower = paragraph.lower()
            
            for keyword in keywords:
                if keyword in paragraph_lower:
                    score += 1
            
            if score > best_score and score >= 2:  # Need at least 2 keyword matches
                best_score = score
                best_paragraph = paragraph.strip()
        
        return best_paragraph
    
    def process_source_files(self, source_files: Dict[str, str]) -> Dict[str, Any]:
        """Enhanced source file processing with better error handling and validation"""
        all_documents = {}
        combined_source_text = ""
        total_files_processed = 0
        total_content_length = 0
        processing_errors = []
        
        print(f"    📁 Processing {len(source_files)} source files (Ground Truth):")
        
        for file_type, file_path in source_files.items():
            print(f"      📄 {file_type}: {os.path.basename(file_path)}")
            
            if not os.path.exists(file_path):
                error_msg = f"File not found: {file_path}"
                print(f"         ❌ {error_msg}")
                processing_errors.append(f"{file_type}: {error_msg}")
                continue
            
            # Check file size and accessibility
            try:
                file_size = os.path.getsize(file_path)
                if file_size == 0:
                    error_msg = f"File is empty: {file_path}"
                    print(f"         ❌ {error_msg}")
                    processing_errors.append(f"{file_type}: {error_msg}")
                    continue
                
                # Check if file is too large (>100MB)
                if file_size > 100 * 1024 * 1024:
                    error_msg = f"File too large ({file_size / (1024*1024):.1f} MB): {file_path}"
                    print(f"         ⚠️ {error_msg}")
                    processing_errors.append(f"{file_type}: {error_msg}")
                    continue
                
                print(f"         📊 File size: {file_size / (1024*1024):.1f} MB")
                
            except Exception as e:
                error_msg = f"Cannot access file: {e}"
                print(f"         ❌ {error_msg}")
                processing_errors.append(f"{file_type}: {error_msg}")
                continue
            
            # Process based on file type - ENHANCED: Pass source file type
            doc_data = None
            try:
                if file_path.lower().endswith('.pdf'):
                    doc_data = self.extract_pdf_text(file_path, file_type)
                elif file_path.lower().endswith(('.pptx', '.ppt')):
                    doc_data = self.extract_pptx_text(file_path, file_type)
                else:
                    error_msg = f"Unsupported file type: {os.path.splitext(file_path)[1]}"
                    print(f"         ❌ {error_msg}")
                    processing_errors.append(f"{file_type}: {error_msg}")
                    continue
            except Exception as e:
                error_msg = f"Processing error: {e}"
                print(f"         ❌ {error_msg}")
                processing_errors.append(f"{file_type}: {error_msg}")
                continue
            
            # Check if extraction was successful
            if doc_data and doc_data.get('extraction_success', False) and doc_data.get('combined_text', '').strip():
                # Store individual document data
                all_documents[file_type] = doc_data
                total_files_processed += 1
                total_content_length += len(doc_data['combined_text'])
                
                # Combine all source text with enhanced separators
                combined_source_text += f"\n\n{'='*80}\n"
                combined_source_text += f"SOURCE DOCUMENT: {file_type.upper().replace('_', ' ')}\n"
                combined_source_text += f"FILE: {os.path.basename(file_path)}\n"
                combined_source_text += f"TYPE: {doc_data.get('pages', doc_data.get('slides', 'unknown'))} {'pages' if 'pages' in doc_data else 'slides'}\n"
                combined_source_text += f"LENGTH: {len(doc_data['combined_text']):,} characters\n"
                combined_source_text += f"EXTRACTION_STATUS: SUCCESS\n"
                combined_source_text += f"{'='*80}\n\n"
                combined_source_text += doc_data['combined_text']
                
                print(f"         ✅ Successfully processed - {len(doc_data['combined_text']):,} characters")
                
                # Validate key information extraction
                key_info = self.extract_key_information(doc_data['combined_text'])
                total_keys = sum(len(v) for v in key_info.values())
                print(f"         🔍 Extracted {total_keys} key information items")
                
            else:
                error_msg = doc_data.get('error', 'No content extracted from file')
                print(f"         ❌ {error_msg}")
                processing_errors.append(f"{file_type}: {error_msg}")
        
        # Enhanced summary with validation
        print(f"    📊 Processing Summary:")
        print(f"       ✅ Files processed: {total_files_processed}/{len(source_files)}")
        print(f"       📄 Total content: {total_content_length:,} characters")
        print(f"       🔍 Combined text length: {len(combined_source_text):,} characters")
        
        if processing_errors:
            print(f"       ⚠️ Errors encountered: {len(processing_errors)}")
            for error in processing_errors:
                print(f"          - {error}")
        
        # Additional validation
        validation_passed = combined_source_text and len(combined_source_text) > 1000
        if validation_passed:
            print(f"       ✅ Source text validation: PASSED")
        else:
            print(f"       ⚠️ Source text validation: WARNING - Combined text may be insufficient")
        
        return {
            "individual_documents": all_documents,
            "combined_source_text": combined_source_text,
            "files_processed": total_files_processed,
            "total_content_length": total_content_length,
            "processing_errors": processing_errors,
            "success_rate": (total_files_processed / len(source_files)) * 100 if source_files else 0,
            "validation_passed": validation_passed
        }
    
    def process_ai_generated_memo(self, memo_path: str) -> Dict[str, Any]:
        """Enhanced AI memo processing with better validation and error handling"""
        print(f"    🤖 Processing AI-generated memo: {os.path.basename(memo_path)}")
        
        if not os.path.exists(memo_path):
            error_msg = f"AI memo not found: {memo_path}"
            print(f"         ❌ {error_msg}")
            return {
                "memo_text": "",
                "memo_sections": {},
                "file_path": memo_path,
                "sections_found": 0,
                "total_length": 0,
                "extraction_success": False,
                "error": error_msg
            }
        
        # Check file size and accessibility
        try:
            file_size = os.path.getsize(memo_path)
            if file_size == 0:
                error_msg = "AI memo file is empty"
                print(f"         ❌ {error_msg}")
                return {
                    "memo_text": "",
                    "memo_sections": {},
                    "file_path": memo_path,
                    "sections_found": 0,
                    "total_length": 0,
                    "extraction_success": False,
                    "error": error_msg
                }
            
            # Check if file is too large
            if file_size > 50 * 1024 * 1024:  # 50MB limit for memos
                error_msg = f"AI memo file too large ({file_size / (1024*1024):.1f} MB)"
                print(f"         ⚠️ {error_msg}")
                return {
                    "memo_text": "",
                    "memo_sections": {},
                    "file_path": memo_path,
                    "sections_found": 0,
                    "total_length": 0,
                    "extraction_success": False,
                    "error": error_msg
                }
            
            print(f"         📊 Memo file size: {file_size / (1024*1024):.1f} MB")
            
        except Exception as e:
            error_msg = f"Cannot access memo file: {e}"
            print(f"         ❌ {error_msg}")
            return {
                "memo_text": "",
                "memo_sections": {},
                "file_path": memo_path,
                "sections_found": 0,
                "total_length": 0,
                "extraction_success": False,
                "error": error_msg
            }
        
        # Extract text from AI-generated memo
        try:
            memo_data = self.extract_pdf_text(memo_path, "ai_memo")
            
            if not memo_data.get('extraction_success', False):
                error_msg = f"Failed to extract text from memo: {memo_data.get('error', 'Unknown error')}"
                print(f"         ❌ {error_msg}")
                return {
                    "memo_text": "",
                    "memo_sections": {},
                    "file_path": memo_path,
                    "sections_found": 0,
                    "total_length": 0,
                    "extraction_success": False,
                    "error": error_msg
                }
            
            memo_text = memo_data['combined_text']
            
            if not memo_text.strip():
                error_msg = "No text content extracted from AI memo"
                print(f"         ❌ {error_msg}")
                return {
                    "memo_text": "",
                    "memo_sections": {},
                    "file_path": memo_path,
                    "sections_found": 0,
                    "total_length": 0,
                    "extraction_success": False,
                    "error": error_msg
                }
            
            # Extract sections from the AI-generated memo using IMPROVED method
            print(f"         🔍 IMPROVED: Extracting sections from {len(memo_text):,} characters...")
            memo_sections = self.extract_memo_sections_from_ai_output(memo_text)
            
            # Count non-empty sections
            sections_found = len([s for s in memo_sections.values() if s.strip()])
            
            print(f"         ✅ IMPROVED memo processing complete:")
            print(f"            📄 Text extracted: {len(memo_text):,} characters")
            print(f"            📋 Sections found: {sections_found}/{len(memo_sections)}")
            
            return {
                "memo_text": memo_text,
                "memo_sections": memo_sections,
                "file_path": memo_path,
                "sections_found": sections_found,
                "total_length": len(memo_text),
                "pages": memo_data.get('pages', 0),
                "extraction_success": True
            }
            
        except Exception as e:
            error_msg = f"Error processing memo: {e}"
            print(f"         ❌ {error_msg}")
            import traceback
            traceback.print_exc()
            return {
                "memo_text": "",
                "memo_sections": {},
                "file_path": memo_path,
                "sections_found": 0,
                "total_length": 0,
                "extraction_success": False,
                "error": error_msg
            }
    
    def process_company_documents(self, company_config: Dict) -> Dict[str, Any]:
        """Enhanced company document processing with comprehensive validation and debugging"""
        company_name = company_config['name']
        result = {
            "company_name": company_name,
            "source_of_truth": {},
            "ai_generated_memo": {},
            "combined_source_text": "",
            "processing_summary": {}
        }
        
        print(f"  📁 Document Processing for {company_name}:")
        
        # Process SOURCE FILES (Source of Truth)
        if 'source_files' in company_config and company_config['source_files']:
            print(f"     🎯 Processing {len(company_config['source_files'])} source files...")
            try:
                source_data = self.process_source_files(company_config['source_files'])
                result["source_of_truth"] = source_data
                result["combined_source_text"] = source_data["combined_source_text"]
                
                # Enhanced debugging for source text
                if result["combined_source_text"]:
                    sample_text = result["combined_source_text"][:500].replace('\n', ' ')
                    print(f"     🔍 Source text sample: {sample_text}...")
                else:
                    print(f"     ⚠️ WARNING: No combined source text generated!")
            except Exception as e:
                print(f"     ❌ Error processing source files: {e}")
                result["source_of_truth"] = {
                    "individual_documents": {},
                    "combined_source_text": "",
                    "files_processed": 0,
                    "processing_errors": [f"Source processing error: {e}"],
                    "validation_passed": False
                }
        else:
            print(f"     ⚠️ No source files configured for {company_name}")
            print(f"        Expected: 'source_files' key in company config")
        
        # Process AI-GENERATED MEMO (AI Output 2) with IMPROVED extraction
        if 'ai_generated_memo' in company_config and company_config['ai_generated_memo']:
            try:
                memo_data = self.process_ai_generated_memo(company_config['ai_generated_memo'])
                result["ai_generated_memo"] = memo_data
            except Exception as e:
                print(f"     ❌ Error processing AI memo: {e}")
                result["ai_generated_memo"] = {
                    "memo_text": "",
                    "memo_sections": {},
                    "sections_found": 0,
                    "extraction_success": False,
                    "error": f"Memo processing error: {e}"
                }
        else:
            print(f"     ⚠️ No AI-generated memo configured for {company_name}")
            print(f"        Expected: 'ai_generated_memo' key in company config")
        
        # Generate comprehensive processing summary
        try:
            source_files_count = result["source_of_truth"].get("files_processed", 0)
            source_content_length = len(result["combined_source_text"])
            memo_sections_count = result["ai_generated_memo"].get("sections_found", 0)
            memo_content_length = result["ai_generated_memo"].get("total_length", 0)
            
            processing_success = (
                source_files_count > 0 and 
                source_content_length > 1000 and
                result["ai_generated_memo"].get("extraction_success", False)
            )
            
            result["processing_summary"] = {
                "source_files_processed": source_files_count,
                "source_content_length": source_content_length,
                "memo_sections_found": memo_sections_count,
                "memo_content_length": memo_content_length,
                "total_content_length": source_content_length + memo_content_length,
                "processing_success": processing_success,
                "validation_passed": result["source_of_truth"].get("validation_passed", False),
                "source_errors": result["source_of_truth"].get("processing_errors", []),
                "memo_errors": [result["ai_generated_memo"].get("error")] if result["ai_generated_memo"].get("error") else []
            }
            
            print(f"  📊 Processing Summary for {company_name}:")
            print(f"     📁 Source files: {source_files_count} processed")
            print(f"     📄 Source content: {source_content_length:,} characters")
            print(f"     🤖 AI memo sections: {memo_sections_count} found")
            print(f"     📋 Total content: {source_content_length + memo_content_length:,} characters")
            print(f"     ✅ Processing success: {processing_success}")
            
            # Report any errors
            total_errors = len(result["processing_summary"]["source_errors"]) + len(result["processing_summary"]["memo_errors"])
            if total_errors > 0:
                print(f"     ⚠️ Errors encountered: {total_errors}")
                
        except Exception as e:
            print(f"     ❌ Error generating processing summary: {e}")
            result["processing_summary"] = {
                "processing_success": False,
                "error": f"Summary generation error: {e}"
            }
        
        return result
    
    def save_extracted_data(self, company_name: str, data: Dict) -> None:
        """Enhanced data saving with comprehensive debugging information"""
        try:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            
            # Create data directory if it doesn't exist
            os.makedirs("data/extracted", exist_ok=True)
            
            # Save source data summary
            if data.get('source_of_truth'):
                try:
                    source_file = f"data/extracted/{company_name}_source_truth_{timestamp}.json"
                    source_summary = {
                        'company_name': company_name,
                        'extraction_timestamp': timestamp,
                        'files_processed': data['source_of_truth'].get('files_processed', 0),
                        'total_length': len(data.get('combined_source_text', '')),
                        'success_rate': data['source_of_truth'].get('success_rate', 0),
                        'validation_passed': data['source_of_truth'].get('validation_passed', False),
                        'processing_errors': data['source_of_truth'].get('processing_errors', []),
                        'sample_content': data.get('combined_source_text', '')[:1000] + '...' if data.get('combined_source_text') else '',
                        'key_information_sample': self.extract_key_information(data.get('combined_source_text', ''))
                    }
                    
                    with open(source_file, 'w', encoding='utf-8') as f:
                        json.dump(source_summary, f, indent=2, ensure_ascii=False)
                    print(f"     💾 Source data summary saved: {source_file}")
                except Exception as e:
                    print(f"     ⚠️ Could not save source data: {e}")
            
            # Save AI memo data summary
            if data.get('ai_generated_memo'):
                try:
                    memo_file = f"data/extracted/{company_name}_ai_memo_{timestamp}.json"
                    memo_summary = {
                        'company_name': company_name,
                        'extraction_timestamp': timestamp,
                        'sections_found': data['ai_generated_memo'].get('sections_found', 0),
                        'total_length': data['ai_generated_memo'].get('total_length', 0),
                        'extraction_success': data['ai_generated_memo'].get('extraction_success', False),
                        'sections': {k: v[:300] + '...' if len(v) > 300 else v 
                                   for k, v in data['ai_generated_memo'].get('memo_sections', {}).items() if v.strip()},
                        'error': data['ai_generated_memo'].get('error')
                    }
                    
                    with open(memo_file, 'w', encoding='utf-8') as f:
                        json.dump(memo_summary, f, indent=2, ensure_ascii=False)
                    print(f"     💾 AI memo summary saved: {memo_file}")
                except Exception as e:
                    print(f"     ⚠️ Could not save memo data: {e}")
                    
        except Exception as e:
            print(f"     ⚠️ Warning: Could not save extracted data: {e}")